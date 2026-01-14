# python-services/extractor/app.py
"""
SmartClaim File Text Extraction Service
Supports: CSV, Excel, PDF, DOCX, PPT, MD, TXT
Features: OCR for image-based documents using PaddleOCR
"""

from fastapi import FastAPI, File, UploadFile, HTTPException, Form
from fastapi.middleware.cors import CORSMiddleware
from typing import Optional
import pandas as pd
import PyPDF2
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF for better PDF handling and image extraction
import io
import os
import tempfile
import logging
from PIL import Image
import numpy as np

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# ============== SINGLETON OCR MODEL ==============
class OCREngine:
    """Singleton class for PaddleOCR to avoid re-initialization"""
    _instance = None
    _reader = None
    _initialized = False
    
    def __new__(cls):
        if cls._instance is None:
            cls._instance = super().__new__(cls)
        return cls._instance
    
    def get_reader(self, languages: list = ['en', 'fr']):
        """Get or initialize the PaddleOCR reader"""
        if not self._initialized:
            try:
                from paddleocr import PaddleOCR
                logger.info("🔄 Initializing PaddleOCR model (first time only)...")
                # PaddleOCR supports: en, fr, german, korean, japan, chinese, etc.
                # Use 'en' for English, 'fr' for French, 'ch' for Chinese
                lang = 'fr' if 'fr' in languages else 'en'
                self._reader = PaddleOCR(
                    use_angle_cls=True,  # Detect rotated text
                    lang=lang,
                    use_gpu=False,  # CPU mode
                    show_log=False,  # Reduce log spam
                )
                self._initialized = True
                logger.info(f"✅ PaddleOCR initialized successfully (lang={lang})")
            except Exception as e:
                logger.error(f"❌ Failed to initialize PaddleOCR: {e}")
                self._reader = None
        return self._reader
    
    @property
    def is_available(self):
        return self._initialized and self._reader is not None


# Global OCR engine instance
ocr_engine = OCREngine()


app = FastAPI(
    title="SmartClaim File Extractor",
    description="Extract text from various file formats with OCR support",
    version="2.0.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def perform_ocr(image_bytes: bytes) -> str:
    """Perform OCR on image bytes using PaddleOCR"""
    reader = ocr_engine.get_reader()
    if not reader:
        return "[OCR not available]"
    
    try:
        # Convert bytes to numpy array for PaddleOCR
        image = Image.open(io.BytesIO(image_bytes))
        image_np = np.array(image.convert('RGB'))
        
        # Perform OCR - PaddleOCR returns list of [box, (text, confidence)]
        results = reader.ocr(image_np, cls=True)
        
        # Extract text from results
        if results and results[0]:
            text_lines = [line[1][0] for line in results[0] if line[1][1] > 0.5]  # Filter by confidence
            return "\n".join(text_lines)
        return ""
    except Exception as e:
        logger.error(f"OCR failed: {e}")
        return f"[OCR error: {str(e)}]"


def perform_ocr_on_image(image: Image.Image) -> str:
    """Perform OCR on PIL Image using PaddleOCR"""
    reader = ocr_engine.get_reader()
    if not reader:
        return "[OCR not available]"
    
    try:
        image_np = np.array(image.convert('RGB'))
        
        # Perform OCR - PaddleOCR returns list of [box, (text, confidence)]
        results = reader.ocr(image_np, cls=True)
        
        # Extract text from results
        if results and results[0]:
            text_lines = [line[1][0] for line in results[0] if line[1][1] > 0.5]  # Filter by confidence
            return "\n".join(text_lines)
        return ""
    except Exception as e:
        logger.error(f"OCR failed: {e}")
        return f"[OCR error: {str(e)}]"


def extract_from_csv(file_content: bytes, user_context: str = "") -> dict:
    """Extract text from CSV file with analysis context"""
    try:
        df = pd.read_csv(io.BytesIO(file_content))
        
        # Generate table summary
        table_text = df.to_string()
        
        # Add analysis context for LLM
        analysis_prompt = f"""
=== CSV/TABLE DATA ANALYSIS ===
The following table data was extracted from an uploaded CSV file.
Number of rows: {len(df)}
Number of columns: {len(df.columns)}
Columns: {', '.join(df.columns.tolist())}

--- TABLE DATA ---
{table_text}

--- ANALYSIS INSTRUCTIONS ---
Please analyze this data for any potential issues, anomalies, or problems that may be relevant to the user's claim or complaint.
{f"User context: {user_context}" if user_context else ""}
"""
        return {"text": analysis_prompt, "is_tabular": True, "rows": len(df), "cols": len(df.columns)}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"CSV parsing error: {str(e)}")


def extract_from_excel(file_content: bytes, user_context: str = "") -> dict:
    """Extract text from Excel file with all sheets and analysis context"""
    try:
        # Read all sheets
        excel_file = pd.ExcelFile(io.BytesIO(file_content))
        all_sheets_text = []
        total_rows = 0
        total_cols = 0
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            total_rows += len(df)
            total_cols = max(total_cols, len(df.columns))
            
            sheet_text = f"\n=== Sheet: {sheet_name} ===\n"
            sheet_text += f"Columns: {', '.join(df.columns.astype(str).tolist())}\n"
            sheet_text += df.to_string()
            all_sheets_text.append(sheet_text)
        
        combined_text = "\n".join(all_sheets_text)
        
        # Add analysis context for LLM
        analysis_prompt = f"""
=== EXCEL DATA ANALYSIS ===
The following data was extracted from an uploaded Excel file.
Number of sheets: {len(excel_file.sheet_names)}
Sheet names: {', '.join(excel_file.sheet_names)}
Total rows across all sheets: {total_rows}

--- SPREADSHEET DATA ---
{combined_text}

--- ANALYSIS INSTRUCTIONS ---
Please analyze this spreadsheet data carefully for:
1. Any anomalies, inconsistencies, or errors in the data
2. Values that seem out of range or suspicious
3. Missing data or incomplete records
4. Any patterns that might indicate problems
{f"User context for this analysis: {user_context}" if user_context else ""}
"""
        return {"text": analysis_prompt, "is_tabular": True, "rows": total_rows, "sheets": len(excel_file.sheet_names)}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Excel parsing error: {str(e)}")


def extract_from_pdf(file_content: bytes) -> dict:
    """Extract text from PDF using PyMuPDF, with OCR fallback for image-based pages"""
    try:
        doc = fitz.open(stream=file_content, filetype="pdf")
        all_text = []
        ocr_used = False
        page_count = len(doc)  # Store page count before closing
        
        for page_num in range(page_count):
            page = doc[page_num]
            page_text = page.get_text("text").strip()
            
            # Check if page has meaningful text
            if len(page_text) > 50:
                # Text-based page - use extracted text
                all_text.append(f"\n--- Page {page_num + 1} ---\n{page_text}")
            else:
                # Image-based page - use OCR
                logger.info(f"[PDF] Page {page_num + 1} appears to be image-based, using OCR...")
                ocr_used = True
                
                # Render page to image
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better OCR
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                
                ocr_text = perform_ocr_on_image(img)
                if ocr_text:
                    all_text.append(f"\n--- Page {page_num + 1} (OCR) ---\n{ocr_text}")
                else:
                    all_text.append(f"\n--- Page {page_num + 1} ---\n[No text detected]")
        
        doc.close()
        
        final_text = "\n".join(all_text).strip()
        return {"text": final_text, "ocr_used": ocr_used, "pages": page_count}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PDF parsing error: {str(e)}")


def extract_from_docx(file_content: bytes) -> dict:
    """Extract text from DOCX file including tables and images"""
    try:
        doc = Document(io.BytesIO(file_content))
        paragraphs = []
        tables_text = []
        images_text = []
        ocr_used = False
        
        # Extract paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                paragraphs.append(paragraph.text)
        
        # Extract tables
        for table_idx, table in enumerate(doc.tables, 1):
            table_rows = []
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                table_rows.append(" | ".join(row_text))
            tables_text.append(f"\n--- Table {table_idx} ---\n" + "\n".join(table_rows))
        
        # Extract images and OCR them
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_data = rel.target_part.blob
                    ocr_text = perform_ocr(image_data)
                    if ocr_text and ocr_text != "[OCR not available]":
                        images_text.append(f"\n--- Image (OCR) ---\n{ocr_text}")
                        ocr_used = True
                except Exception as e:
                    logger.warning(f"Could not process image in DOCX: {e}")
        
        # Combine all text
        full_text = "\n".join(paragraphs)
        if tables_text:
            full_text += "\n\n=== TABLES ===" + "".join(tables_text)
        if images_text:
            full_text += "\n\n=== IMAGES (OCR) ===" + "".join(images_text)
        
        return {"text": full_text.strip(), "ocr_used": ocr_used}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"DOCX parsing error: {str(e)}")


def extract_from_pptx(file_content: bytes) -> dict:
    """Extract text from PowerPoint including OCR for image slides"""
    try:
        prs = Presentation(io.BytesIO(file_content))
        all_text = []
        ocr_used = False
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text_parts = []
            has_text = False
            
            for shape in slide.shapes:
                # Extract text from text frames
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text)
                    has_text = True
                
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        slide_text_parts.append(" | ".join(row_text))
                    has_text = True
                
                # OCR images in slides
                if hasattr(shape, "image"):
                    try:
                        image_bytes = shape.image.blob
                        ocr_text = perform_ocr(image_bytes)
                        if ocr_text and ocr_text not in ["[OCR not available]", ""]:
                            slide_text_parts.append(f"[Image OCR]: {ocr_text}")
                            ocr_used = True
                    except Exception as e:
                        logger.warning(f"Could not OCR image in slide {slide_num}: {e}")
            
            slide_content = "\n".join(slide_text_parts) if slide_text_parts else "[No text content]"
            all_text.append(f"\n--- Slide {slide_num} ---\n{slide_content}")
        
        return {"text": "\n".join(all_text).strip(), "ocr_used": ocr_used, "slides": len(prs.slides)}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PPTX parsing error: {str(e)}")


def extract_from_text(file_content: bytes) -> dict:
    """Extract text from plain text or markdown files"""
    try:
        text = file_content.decode('utf-8')
        return {"text": text}
    except UnicodeDecodeError:
        try:
            text = file_content.decode('latin-1')
            return {"text": text}
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Text decoding error: {str(e)}")


def extract_from_image(file_content: bytes) -> dict:
    """Extract text from standalone image files using OCR"""
    try:
        ocr_text = perform_ocr(file_content)
        return {"text": ocr_text, "ocr_used": True}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Image OCR error: {str(e)}")


@app.post("/extract")
async def extract_text(
    file: UploadFile = File(...),
    user_context: Optional[str] = Form(default="")
):
    """
    Extract text from uploaded file with intelligent OCR support
    
    Supported formats: CSV, Excel (XLS/XLSX), PDF, DOCX, PPTX, TXT, MD, Images (PNG/JPG/JPEG)
    
    Features:
    - Automatic OCR for image-based documents
    - Multi-sheet Excel support with analysis context
    - Table extraction from Word/PowerPoint
    - Image extraction and OCR from embedded images
    """
    file_content = await file.read()
    file_name = file.filename.lower() if file.filename else ""
    
    logger.info(f"📥 Processing file: {file.filename} ({len(file_content)} bytes)")
    
    result = {"text": "", "ocr_used": False}
    metadata = {
        "file_name": file.filename,
        "file_type": file.content_type,
        "size_bytes": len(file_content)
    }
    
    try:
        # Determine file type and extract accordingly
        if file_name.endswith('.csv'):
            result = extract_from_csv(file_content, user_context)
            metadata["format"] = "csv"
            
        elif file_name.endswith(('.xls', '.xlsx')):
            result = extract_from_excel(file_content, user_context)
            metadata["format"] = "excel"
            
        elif file_name.endswith('.pdf'):
            result = extract_from_pdf(file_content)
            metadata["format"] = "pdf"
            
        elif file_name.endswith(('.doc', '.docx')):
            result = extract_from_docx(file_content)
            metadata["format"] = "docx"
            
        elif file_name.endswith(('.ppt', '.pptx')):
            result = extract_from_pptx(file_content)
            metadata["format"] = "pptx"
            
        elif file_name.endswith(('.txt', '.md')):
            result = extract_from_text(file_content)
            metadata["format"] = "text"
            
        elif file_name.endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.webp')):
            result = extract_from_image(file_content)
            metadata["format"] = "image"
            
        else:
            raise HTTPException(
                status_code=400, 
                detail=f"Unsupported file type: {file_name}. Supported: CSV, Excel, PDF, DOCX, PPTX, TXT, MD, Images"
            )
        
        text = result.get("text", "")
        metadata["extracted_length"] = len(text)
        metadata["word_count"] = len(text.split())
        metadata["ocr_used"] = result.get("ocr_used", False)
        
        # Add any extra metadata from extraction
        for key in ["pages", "slides", "sheets", "rows", "cols", "is_tabular"]:
            if key in result:
                metadata[key] = result[key]
        
        logger.info(f"✅ Extracted {metadata['word_count']} words from {file.filename} (OCR: {metadata['ocr_used']})")
        
        return {
            "success": True,
            "text": text,
            "metadata": metadata
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Extraction error for {file.filename}: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Extraction error: {str(e)}")


@app.post("/extract-multiple")
async def extract_multiple_files(files: list[UploadFile] = File(...)):
    """
    Extract text from multiple files at once
    
    Returns combined text from all files with clear separators
    """
    all_results = []
    combined_text = []
    total_words = 0
    
    for file in files:
        try:
            file_content = await file.read()
            file_name = file.filename.lower() if file.filename else ""
            
            logger.info(f"📥 Processing: {file.filename}")
            
            result = {"text": "", "ocr_used": False}
            
            if file_name.endswith('.csv'):
                result = extract_from_csv(file_content)
            elif file_name.endswith(('.xls', '.xlsx')):
                result = extract_from_excel(file_content)
            elif file_name.endswith('.pdf'):
                result = extract_from_pdf(file_content)
            elif file_name.endswith(('.doc', '.docx')):
                result = extract_from_docx(file_content)
            elif file_name.endswith(('.ppt', '.pptx')):
                result = extract_from_pptx(file_content)
            elif file_name.endswith(('.txt', '.md')):
                result = extract_from_text(file_content)
            elif file_name.endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.webp')):
                result = extract_from_image(file_content)
            else:
                result = {"text": f"[Unsupported file type: {file_name}]", "ocr_used": False}
            
            text = result.get("text", "")
            word_count = len(text.split())
            total_words += word_count
            
            all_results.append({
                "file_name": file.filename,
                "word_count": word_count,
                "ocr_used": result.get("ocr_used", False),
                "success": True
            })
            
            combined_text.append(f"\n\n{'='*60}\n=== FILE: {file.filename} ===\n{'='*60}\n\n{text}")
            
        except Exception as e:
            logger.error(f"Error processing {file.filename}: {e}")
            all_results.append({
                "file_name": file.filename,
                "success": False,
                "error": str(e)
            })
    
    return {
        "success": True,
        "text": "\n".join(combined_text).strip(),
        "files_processed": len(files),
        "total_words": total_words,
        "results": all_results
    }


@app.get("/health")
async def health_check():
    """Health check endpoint"""
    return {
        "status": "healthy", 
        "service": "file-extractor",
        "version": "2.0.0",
        "ocr_available": ocr_engine.is_available,
        "supported_formats": ["csv", "xls", "xlsx", "pdf", "doc", "docx", "ppt", "pptx", "txt", "md", "png", "jpg", "jpeg"]
    }


@app.get("/")
async def root():
    """Root endpoint with service info"""
    return {
        "service": "SmartClaim File Extractor",
        "version": "2.0.0",
        "description": "Extract text from documents with intelligent OCR support",
        "endpoints": {
            "POST /extract": "Extract text from a single file",
            "POST /extract-multiple": "Extract text from multiple files",
            "GET /health": "Health check"
        }
    }


@app.on_event("startup")
async def startup_event():
    """Pre-initialize OCR on startup for faster first request"""
    logger.info("🚀 Starting file extractor service...")
    # Optionally pre-initialize OCR (comment out if you want lazy loading)
    # ocr_engine.get_reader()


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)